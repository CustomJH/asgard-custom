#!/usr/bin/env python3
"""Fill placeholders in an existing .docx / .pptx / .xlsx, formatting intact.

    asgard skills run asgard-office -- fill FORM.docx --values v.json -o out.docx

This is the path for a document somebody else designed — an employer's letterhead,
a client's proposal shell, a government form. The layout is theirs and must not
move; only the marked slots change.

Two mechanics:

  {{field}}                     scalar substitution, anywhere text lives
  {{#rows}} … {{/rows}}         in a table row: repeat that row once per item,
                                with each item's keys in scope

Word fragments a visible phrase across many runs (revision ids, spell-check
state), so `{{client}}` frequently does not exist as one string in the XML. The
substitution here works on the paragraph's joined text and writes back only the
runs the span actually touches — which is why the surrounding bold, colour, and
font survive.
"""

from __future__ import annotations

import argparse
import copy
import json
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from officelib import specs, templates  # noqa: E402

ROW_LOOP = re.compile(r"\{\{\s*#\s*([A-Za-z0-9_.\-]+)\s*\}\}")
ROW_LOOP_END = re.compile(r"\{\{\s*/\s*([A-Za-z0-9_.\-]+)\s*\}\}")
SCALAR = re.compile(r"\{\{\s*([A-Za-z0-9_.\-]+)\s*\}\}")


class Report:
    def __init__(self) -> None:
        self.filled = 0
        self.rows = 0
        self.missing: set[str] = set()
        self.remaining: set[str] = set()

    def as_dict(self, output: str) -> dict:
        return {
            "output": output,
            "substitutions": self.filled,
            "repeated_rows": self.rows,
            "missing_fields": sorted(self.missing),
            "unfilled_placeholders": sorted(self.remaining),
        }


def _resolve(values: dict, name: str):
    cursor = values
    for part in name.split("."):
        if isinstance(cursor, dict) and part in cursor:
            cursor = cursor[part]
        else:
            return None
    return cursor


def fill_paragraph(runs, values: dict, report: Report) -> None:
    """Substitute across run boundaries without disturbing runs the span misses."""
    texts = [run.text or "" for run in runs]
    joined = "".join(texts)
    if "{{" not in joined:
        return
    matches = list(SCALAR.finditer(joined))
    if not matches:
        report.remaining.update(match.group(0) for match in ROW_LOOP.finditer(joined))
        return
    # Map every character position to (run index, offset within the run).
    bounds: list[tuple[int, int]] = []
    for index, text in enumerate(texts):
        bounds.append((sum(len(item) for item in texts[:index]), index))
    replacements: list[tuple[int, int, str]] = []
    for match in matches:
        name = match.group(1)
        value = _resolve(values, name)
        if value is None:
            report.missing.add(name)
            report.remaining.add(match.group(0))
            continue
        report.filled += 1
        replacements.append((match.start(), match.end(), templates.stringify(value)))
    if not replacements:
        return
    # Rewrite right to left so earlier offsets stay valid.
    buffers = list(texts)
    for start, end, text in reversed(replacements):
        start_run = _run_at(bounds, texts, start)
        end_run = _run_at(bounds, texts, end - 1)
        start_offset = start - bounds[start_run][0]
        end_offset = end - bounds[end_run][0]
        if start_run == end_run:
            buffers[start_run] = buffers[start_run][:start_offset] + text + buffers[start_run][end_offset:]
        else:
            buffers[start_run] = buffers[start_run][:start_offset] + text
            for middle in range(start_run + 1, end_run):
                buffers[middle] = ""
            buffers[end_run] = buffers[end_run][end_offset:]
        texts = buffers
        bounds = [(sum(len(item) for item in texts[:index]), index) for index in range(len(texts))]
    for run, text in zip(runs, buffers):
        if run.text != text:
            run.text = text


def _run_at(bounds, texts, position: int) -> int:
    for index in range(len(texts) - 1, -1, -1):
        if bounds[index][0] <= position:
            return index
    return 0


# ---------------------------------------------------------------------- docx


def _row_loop_name(row) -> str:
    text = " ".join(cell.text for cell in row.cells)
    match = ROW_LOOP.search(text)
    return match.group(1) if match and ROW_LOOP_END.search(text) else ""


def _strip_markers(row) -> None:
    for cell in row.cells:
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    run.text = ROW_LOOP_END.sub("", ROW_LOOP.sub("", run.text))


def fill_docx(path: Path, values: dict, out: Path, report: Report) -> None:
    from docx import Document

    document = Document(str(path))
    for table in document.tables:
        _expand_docx_rows(table, values, report)
    for paragraph in document.paragraphs:
        fill_paragraph(paragraph.runs, values, report)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    fill_paragraph(paragraph.runs, values, report)
    for section in document.sections:
        for container in (section.header, section.footer, section.first_page_header, section.first_page_footer):
            if container is None:
                continue
            for paragraph in container.paragraphs:
                fill_paragraph(paragraph.runs, values, report)
    out.parent.mkdir(parents=True, exist_ok=True)
    document.save(str(out))


def _expand_docx_rows(table, values: dict, report: Report) -> None:
    for row in list(table.rows):
        name = _row_loop_name(row)
        if not name:
            continue
        items = _resolve(values, name)
        if not isinstance(items, list):
            report.missing.add(name)
            _strip_markers(row)
            continue
        anchor = row._tr
        for item in items:
            clone = copy.deepcopy(anchor)
            anchor.addprevious(clone)
            report.rows += 1
        anchor.getparent().remove(anchor)
        # Re-read the table so the clones are addressable, then fill each in turn.
        fresh = [candidate for candidate in table.rows if ROW_LOOP.search(" ".join(c.text for c in candidate.cells))]
        for candidate, item in zip(fresh, items):
            _strip_markers(candidate)
            scope = {**values, **item} if isinstance(item, dict) else {**values, "item": item}
            for cell in candidate.cells:
                for paragraph in cell.paragraphs:
                    fill_paragraph(paragraph.runs, scope, report)


# ---------------------------------------------------------------------- pptx


def fill_pptx(path: Path, values: dict, out: Path, report: Report) -> None:
    from pptx import Presentation

    presentation = Presentation(str(path))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    fill_paragraph(paragraph.runs, values, report)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            fill_paragraph(paragraph.runs, values, report)
        if slide.has_notes_slide:
            for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                fill_paragraph(paragraph.runs, values, report)
    out.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(out))


# ---------------------------------------------------------------------- xlsx


def fill_xlsx(path: Path, values: dict, out: Path, report: Report) -> None:
    import openpyxl

    workbook = openpyxl.load_workbook(str(path), keep_vba=path.suffix.lower() == ".xlsm")
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if not isinstance(cell.value, str) or "{{" not in cell.value:
                    continue
                text = cell.value
                for match in SCALAR.finditer(cell.value):
                    resolved = _resolve(values, match.group(1))
                    if resolved is None:
                        report.missing.add(match.group(1))
                        report.remaining.add(match.group(0))
                        continue
                    report.filled += 1
                    text = text.replace(match.group(0), templates.stringify(resolved))
                cell.value = text
    out.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(str(out))


def fill(path: Path, values: dict, out: Path) -> dict:
    report = Report()
    suffix = path.suffix.lower()
    if suffix in (".docx", ".dotx"):
        fill_docx(path, values, out, report)
    elif suffix in (".pptx", ".potx"):
        fill_pptx(path, values, out, report)
    elif suffix in (".xlsx", ".xlsm"):
        fill_xlsx(path, values, out, report)
    else:
        raise ValueError(f"fill supports docx, pptx, and xlsx; got {suffix}")
    return report.as_dict(str(out))


def scan_typed(path: Path) -> list[tuple[str, str, list[str]]]:
    """Placeholders a form declares, as (name, type, row keys).

    Reads the package XML rather than the rendered text: a placeholder split
    across runs is still one `{{name}}` once the tags are stripped, and header,
    footer, and notes parts never reach a text extractor.

    A name used as `{{#rows}} … {{/rows}}` is a list, and the names *inside*
    that span belong to each row rather than to the document — reporting them as
    top-level fields is how an adopted template ends up with a schema nobody can
    satisfy.
    """
    import zipfile
    from xml.etree import ElementTree

    scalars: list[str] = []
    sections: dict[str, list[str]] = {}
    with zipfile.ZipFile(path) as archive:
        for entry in sorted(archive.namelist()):
            if not entry.endswith(".xml") or entry.startswith("docProps/"):
                continue
            try:
                root = ElementTree.fromstring(archive.read(entry))
            except ElementTree.ParseError:
                continue
            text = "".join(node.text or "" for node in root.iter() if node.text)
            inner_spans: list[tuple[int, int, str]] = []
            for opener in ROW_LOOP.finditer(text):
                closer = ROW_LOOP_END.search(text, opener.end())
                while closer is not None and closer.group(1) != opener.group(1):
                    closer = ROW_LOOP_END.search(text, closer.end())
                if closer is not None:
                    inner_spans.append((opener.end(), closer.start(), opener.group(1)))
                    sections.setdefault(opener.group(1), [])
            for match in SCALAR.finditer(text):
                name = match.group(1)
                owner = next((label for start, end, label in inner_spans if start <= match.start() < end), "")
                if owner:
                    if name not in sections[owner]:
                        sections[owner].append(name)
                elif name not in scalars and name not in sections:
                    scalars.append(name)
    return [(name, "text", []) for name in scalars if name not in sections] + [
        (name, "table", keys) for name, keys in sections.items()
    ]


def scan(path: Path) -> list[str]:
    """Flat placeholder names, row keys included — the quick answer for `--scan`."""
    out: list[str] = []
    for name, _, keys in scan_typed(path):
        out.append(name)
        out.extend(key for key in keys if key not in out)
    return out


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(prog="fill", description="Fill {{placeholders}} in an existing office file")
    parser.add_argument("file", type=Path)
    parser.add_argument("-o", "--output", type=Path)
    parser.add_argument("--values", type=Path, help="JSON/YAML values file")
    parser.add_argument("--scan", action="store_true", help="list the placeholders instead of filling them")
    parser.add_argument("--strict", action="store_true", help="fail when any placeholder is left unfilled")
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args(argv)
    try:
        if args.scan:
            names = scan(args.file)
            print(json.dumps(names, ensure_ascii=False, indent=2) if args.json else "\n".join(names) or "(none)")
            return 0
        if not args.values or not args.output:
            print("fill needs --values and --output (or --scan)", file=sys.stderr)
            return 2
        report = fill(args.file, specs.load_data(args.values), args.output)
    except (ValueError, OSError, KeyError, json.JSONDecodeError) as exc:
        print(f"fill failed: {exc}", file=sys.stderr)
        return 1
    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        print(f"wrote {report['output']}  ({report['substitutions']} substitutions, {report['repeated_rows']} rows)")
        if report["missing_fields"]:
            print("  no value supplied for: " + ", ".join(report["missing_fields"]))
    return 1 if args.strict and report["unfilled_placeholders"] else 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
