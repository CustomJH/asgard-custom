#!/usr/bin/env python3
"""Markdown deck spec -> .pptx.

    asgard skills run asgard-office -- build pptx DECK.md -o out.pptx [--template NAME] [--values v.json]

Slides are separated by `---`, the convention every Markdown deck tool already
uses. Per-slide direction rides in HTML comments so the spec still reads as a
document:

    ## Cost overrun
    <!-- layout: stat -->
    <!-- notes: lead with the number, not the cause -->
    - 107% :: of plan
    - 3.2x :: retry amplification

Every slide is drawn on a blank layout with explicit geometry. Nothing inherits
from a master unless a template supplies one, so what the spec says is what the
deck is.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from officelib import color, compose, mdblocks, pptxkit, specs  # noqa: E402

LAYOUTS = ("title", "section", "bullets", "two-col", "stat", "quote", "table", "image", "blank")
MARGIN_IN = 0.6
COLUMN_SPLIT = "|||"
_DIRECTIVE = re.compile(r"^\s*(layout|notes|background|image|columns)\s*:\s*(.*)$", re.I | re.S)
_STAT = re.compile(r"^(.*?)\s*::\s*(.*)$")


@dataclass
class Slide:
    title: list[mdblocks.Run] = field(default_factory=list)
    blocks: list[mdblocks.Block] = field(default_factory=list)
    layout: str = ""
    notes: str = ""
    background: str = ""
    image: str = ""


def split_slides(body: str) -> list[Slide]:
    """`---` at column zero starts a new slide; everything else is block content."""
    slides: list[Slide] = []
    current = Slide()
    for block in mdblocks.parse(body):
        if block.kind == "rule":
            if current.title or current.blocks or current.layout:
                slides.append(current)
            current = Slide()
            continue
        if block.kind == "comment":
            match = _DIRECTIVE.match(block.text)
            if match:
                key, value = match.group(1).lower(), match.group(2).strip()
                if key == "layout":
                    current.layout = value.lower()
                elif key == "notes":
                    current.notes = (current.notes + "\n" + value).strip()
                elif key == "background":
                    current.background = specs.hex_color(value, "")
                elif key == "image":
                    current.image = value
            continue
        if block.kind == "heading" and not current.title and not current.blocks:
            current.title = block.runs
            continue
        current.blocks.append(block)
    if current.title or current.blocks:
        slides.append(current)
    return slides


def infer_layout(slide: Slide, index: int) -> str:
    if slide.layout in LAYOUTS:
        return slide.layout
    if slide.layout:
        raise ValueError(f"unknown slide layout {slide.layout!r} (expected one of {', '.join(LAYOUTS)})")
    kinds = [block.kind for block in slide.blocks]
    if slide.image or "image" in kinds:
        return "image"
    if "table" in kinds:
        return "table"
    if not slide.blocks:
        return "section" if index else "title"
    if index == 0 and kinds == ["para"]:
        return "title"
    if kinds == ["quote"]:
        return "quote"
    if any(COLUMN_SPLIT in mdblocks.plain(item.runs) for block in slide.blocks for item in block.items):
        return "two-col"
    if "list" in kinds and all(_STAT.match(mdblocks.plain(item.runs)) for block in slide.blocks
                               if block.kind == "list" for item in block.items):
        return "stat"
    return "bullets"


class DeckBuilder:
    def __init__(self, meta: specs.DocMeta, size: tuple[float, float], base: Path | None, assets: Path):
        self.meta = meta
        self.theme = meta.theme
        self.assets = assets
        self.warnings: list[str] = []
        self._grounds: list[str] = []
        self.presentation = Presentation(str(base)) if base else Presentation()
        self.presentation.slide_width = Inches(size[0])
        self.presentation.slide_height = Inches(size[1])
        self.width, self.height = size
        self.blank = self.presentation.slide_layouts[6 if len(self.presentation.slide_layouts) > 6 else 0]

    # ------------------------------------------------------------------ atoms

    @property
    def content_width(self) -> float:
        return self.width - 2 * MARGIN_IN

    def _new(self, background: str) -> object:
        slide = self.presentation.slides.add_slide(self.blank)
        pptxkit.set_background(slide, background)
        self._grounds.append(background)
        return slide

    def _box(self, slide, left: float, top: float, width: float, height: float):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        pptxkit.zero_inset(frame)
        return box

    def _write(self, frame, runs, *, size: float, color: str, bold: bool = False, align=PP_ALIGN.LEFT, first=False):
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        paragraph.alignment = align
        pptxkit.set_no_bullet(paragraph)
        for item in runs:
            if not item.text:
                continue
            run = paragraph.add_run()
            run.text = item.text
            run.font.size = Pt(size)
            run.font.bold = bold or item.bold
            run.font.italic = item.italic
            run.font.name = self.theme.font_mono if item.code else self.theme.font_body
            run.font.color.rgb = RGBColor.from_string(color)
            if item.link:
                run.hyperlink.address = item.link
        return paragraph

    def _title(self, slide, runs, *, color: str, top: float = MARGIN_IN) -> float:
        if not runs:
            return top
        size = self.theme.size_h1 + 18  # 38pt at the default scale
        text = mdblocks.plain(runs)
        lines = pptxkit.estimate_lines(text, font_pt=size, box_width_in=self.content_width)
        height = lines * size * 1.25 / 72.0
        box = self._box(slide, MARGIN_IN, top, self.content_width, height)
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        self._write(box.text_frame, runs, size=size, color=color, bold=True, first=True)
        return top + height + 0.28

    # ---------------------------------------------------------------- layouts

    def title_slide(self, slide_spec: Slide) -> None:
        background = slide_spec.background or self.theme.primary
        slide = self._new(background)
        on_dark = color.is_dark(background)
        ink = self.theme.background if on_dark else self.theme.primary
        muted = self.theme.surface if on_dark else self.theme.muted
        box = self._box(slide, MARGIN_IN, self.height * 0.34, self.content_width, self.height * 0.4)
        runs = slide_spec.title or [mdblocks.Run(self.meta.title)]
        self._write(box.text_frame, runs, size=self.theme.size_title + 4, color=ink, bold=True, first=True)
        for block in slide_spec.blocks:
            if block.kind == "para":
                self._write(box.text_frame, block.runs, size=self.theme.size_h2 + 2, color=muted)
        byline = " · ".join(part for part in (self.meta.author, self.meta.company, self.meta.date) if part)
        if byline:
            footer = self._box(slide, MARGIN_IN, self.height - MARGIN_IN - 0.35, self.content_width, 0.35)
            self._write(footer.text_frame, [mdblocks.Run(byline)], size=12, color=muted, first=True)
        self._notes(slide, slide_spec)

    def section_slide(self, slide_spec: Slide) -> None:
        background = slide_spec.background or self.theme.secondary
        slide = self._new(background)
        ink = color.readable_ink(background, self.theme.background, self.theme.primary)
        box = self._box(slide, MARGIN_IN, self.height * 0.4, self.content_width, 1.6)
        self._write(box.text_frame, slide_spec.title, size=self.theme.size_title, color=ink, bold=True, first=True)
        for block in slide_spec.blocks:
            if block.kind == "para":
                self._write(box.text_frame, block.runs, size=self.theme.size_h2, color=ink)
        self._notes(slide, slide_spec)

    def bullets_slide(self, slide_spec: Slide, *, blocks=None, left=MARGIN_IN, width=None, slide=None) -> None:
        own = slide is None
        if own:
            slide = self._new(slide_spec.background or self.theme.background)
        width = width if width is not None else self.content_width
        top = self._title(slide, slide_spec.title, color=self.theme.primary) if own else MARGIN_IN + 1.15
        box = self._box(slide, left, top, width, self.height - top - MARGIN_IN)
        frame = box.text_frame
        first = True
        for block in blocks if blocks is not None else slide_spec.blocks:
            if block.kind == "list":
                for item in mdblocks.normalise_levels(block.items):
                    size = self.theme.size_body + (5 if item.level == 0 else 3)
                    paragraph = self._write(
                        frame, item.runs, size=size, color=self.theme.text, first=first
                    )
                    first = False
                    paragraph.level = min(item.level, 4)
                    if item.ordered:
                        pptxkit.set_auto_number(paragraph)
                    else:
                        pptxkit.set_bullet(paragraph, color=self.theme.accent)
                    pptxkit.set_indent(paragraph, min(item.level, 4))
                    pptxkit.set_space(paragraph, after=8 if item.level == 0 else 4)
            elif block.kind in ("para", "quote"):
                paragraph = self._write(
                    frame, block.runs, size=self.theme.size_body + 4, color=self.theme.secondary, first=first
                )
                first = False
                pptxkit.set_space(paragraph, after=10)
            elif block.kind == "code":
                for line in block.text.split("\n"):
                    paragraph = self._write(
                        frame, [mdblocks.Run(line or " ", code=True)], size=self.theme.size_body + 1,
                        color=self.theme.secondary, first=first
                    )
                    first = False
                    pptxkit.set_space(paragraph, after=0)
        if own:
            self._notes(slide, slide_spec)

    def two_col_slide(self, slide_spec: Slide) -> None:
        slide = self._new(slide_spec.background or self.theme.background)
        self._title(slide, slide_spec.title, color=self.theme.primary)
        left_blocks, right_blocks = _split_columns(slide_spec.blocks)
        gutter = 0.5
        column = (self.content_width - gutter) / 2
        self.bullets_slide(slide_spec, blocks=left_blocks, left=MARGIN_IN, width=column, slide=slide)
        self.bullets_slide(
            slide_spec, blocks=right_blocks, left=MARGIN_IN + column + gutter, width=column, slide=slide
        )
        self._notes(slide, slide_spec)

    def stat_slide(self, slide_spec: Slide) -> None:
        slide = self._new(slide_spec.background or self.theme.background)
        top = self._title(slide, slide_spec.title, color=self.theme.primary)
        pairs: list[tuple[str, str]] = []
        for block in slide_spec.blocks:
            for item in block.items:
                match = _STAT.match(mdblocks.plain(item.runs))
                if match:
                    pairs.append((match.group(1).strip(), match.group(2).strip()))
        if not pairs:
            self.bullets_slide(slide_spec)
            return
        columns = min(len(pairs), 4)
        gutter = 0.4
        cell = (self.content_width - gutter * (columns - 1)) / columns
        for index, (value, label) in enumerate(pairs[:8]):
            row, column = divmod(index, columns)
            left = MARGIN_IN + column * (cell + gutter)
            offset = top + row * 2.1
            ground = slide_spec.background or self.theme.background
            value_box = self._box(slide, left, offset, cell, 1.1)
            self._write(
                value_box.text_frame, [mdblocks.Run(value)], size=self.theme.size_title + 8,
                color=color.readable_accent(self.theme.accent, ground), bold=True, first=True
            )
            label_box = self._box(slide, left, offset + 1.22, cell, 0.8)
            self._write(
                label_box.text_frame, [mdblocks.Run(label)], size=self.theme.size_body + 2,
                color=color.readable_accent(self.theme.muted, ground, floor=4.5), first=True
            )
        self._notes(slide, slide_spec)

    def quote_slide(self, slide_spec: Slide) -> None:
        background = slide_spec.background or self.theme.surface
        slide = self._new(background)
        ink = color.readable_ink(background, self.theme.background, self.theme.primary)
        top = self._title(slide, slide_spec.title, color=ink) if slide_spec.title else self.height * 0.3
        box = self._box(slide, MARGIN_IN + 0.4, max(top, self.height * 0.28), self.content_width - 0.8, 2.5)
        for block in slide_spec.blocks:
            runs = [mdblocks.Run(mdblocks.plain(block.runs), italic=True)]
            self._write(box.text_frame, runs, size=self.theme.size_h1 + 6, color=ink, first=True)
            break
        self._notes(slide, slide_spec)

    def table_slide(self, slide_spec: Slide) -> None:
        slide = self._new(slide_spec.background or self.theme.background)
        top = self._title(slide, slide_spec.title, color=self.theme.primary)
        source = next((block for block in slide_spec.blocks if block.kind == "table"), None)
        if source is None:
            self.bullets_slide(slide_spec)
            return
        columns = max(len(source.header), max((len(row) for row in source.rows), default=0))
        rows = len(source.rows) + 1
        height = min(self.height - top - MARGIN_IN, 0.45 * rows)
        shape = slide.shapes.add_table(
            rows, columns, Inches(MARGIN_IN), Inches(top), Inches(self.content_width), Inches(height)
        )
        table = shape.table
        for index in range(columns):
            cell = table.cell(0, index)
            cell.text = mdblocks.plain(source.header[index]) if index < len(source.header) else ""
            _style_cell(cell, self.theme, size=self.theme.size_body + 2, bold=True, color=self.theme.primary)
        for row_index, row in enumerate(source.rows, start=1):
            for index in range(columns):
                cell = table.cell(row_index, index)
                cell.text = mdblocks.plain(row[index]) if index < len(row) else ""
                _style_cell(cell, self.theme, size=self.theme.size_body + 1, color=self.theme.text)
        self._notes(slide, slide_spec)

    def image_slide(self, slide_spec: Slide) -> None:
        slide = self._new(slide_spec.background or self.theme.background)
        source = slide_spec.image or next(
            (block.text for block in slide_spec.blocks if block.kind == "image"), ""
        )
        path = Path(source)
        if source and not path.is_absolute():
            path = self.assets / source
        text_blocks = [block for block in slide_spec.blocks if block.kind != "image"]
        half = self.content_width / 2 - 0.2
        top = self._title(slide, slide_spec.title, color=self.theme.primary)
        if path.is_file():
            try:
                slide.shapes.add_picture(
                    str(path), Inches(MARGIN_IN + half + 0.4), Inches(top), width=Inches(half)
                )
            except Exception as exc:
                self.warnings.append(f"image could not be embedded ({exc.__class__.__name__}): {source}")
        elif source:
            self.warnings.append(f"image not found: {source}")
        if text_blocks:
            self.bullets_slide(slide_spec, blocks=text_blocks, left=MARGIN_IN, width=half, slide=slide)
        self._notes(slide, slide_spec)

    def blank_slide(self, slide_spec: Slide) -> None:
        slide = self._new(slide_spec.background or self.theme.background)
        self._title(slide, slide_spec.title, color=self.theme.primary)
        self._notes(slide, slide_spec)

    def _notes(self, slide, slide_spec: Slide) -> None:
        if slide_spec.notes:
            slide.notes_slide.notes_text_frame.text = slide_spec.notes

    def footers(self) -> None:
        if not self.meta.footer:
            return
        for index, slide in enumerate(self.presentation.slides):
            if index == 0:
                continue
            ground = self._grounds[index]
            box = self._box(slide, MARGIN_IN, self.height - 0.45, self.content_width, 0.3)
            frame = box.text_frame
            frame.vertical_anchor = MSO_ANCHOR.BOTTOM
            self._write(
                frame,
                [mdblocks.Run(f"{self.meta.footer}   {index + 1}")],
                size=10,
                color=color.readable_accent(self.theme.muted, ground, floor=4.5),
                align=PP_ALIGN.RIGHT,
                first=True,
            )

    def build(self, slides: list[Slide]) -> None:
        emit = {
            "title": self.title_slide,
            "section": self.section_slide,
            "bullets": self.bullets_slide,
            "two-col": self.two_col_slide,
            "stat": self.stat_slide,
            "quote": self.quote_slide,
            "table": self.table_slide,
            "image": self.image_slide,
            "blank": self.blank_slide,
        }
        for index, slide_spec in enumerate(slides):
            emit[infer_layout(slide_spec, index)](slide_spec)
        self.footers()

    def finish(self, out: Path) -> None:
        core = self.presentation.core_properties
        core.title = self.meta.title or core.title
        core.author = self.meta.author or core.author
        out.parent.mkdir(parents=True, exist_ok=True)
        self.presentation.save(str(out))


def _split_columns(blocks: list[mdblocks.Block]) -> tuple[list, list]:
    """Split on a `|||` marker inside a list, or halfway when there is none."""
    left: list[mdblocks.Block] = []
    right: list[mdblocks.Block] = []
    target = left
    for block in blocks:
        if block.kind != "list":
            target.append(block)
            continue
        bucket_left, bucket_right, seen = [], [], False
        for item in block.items:
            if mdblocks.plain(item.runs).strip() == COLUMN_SPLIT:
                seen = True
                continue
            (bucket_right if seen else bucket_left).append(item)
        if not seen:
            middle = (len(bucket_left) + 1) // 2
            bucket_left, bucket_right = bucket_left[:middle], bucket_left[middle:]
        if bucket_left:
            left.append(mdblocks.Block("list", items=bucket_left, ordered=block.ordered))
        if bucket_right:
            right.append(mdblocks.Block("list", items=bucket_right, ordered=block.ordered))
        target = right
    return left, right


def _style_cell(cell, theme: specs.Theme, *, size: float, color: str, bold: bool = False) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor.from_string(theme.surface if bold else theme.background)
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.name = theme.font_body
            run.font.color.rgb = RGBColor.from_string(color)


def build(
    spec_path: Path,
    out: Path,
    *,
    template_name: str = "",
    values_path: Path | None = None,
    root: Path | None = None,
) -> dict:
    front, body, template, missing = compose.resolve_spec(
        spec_path, template_name, values_path, root, ("pptx", "md")
    )
    meta = specs.DocMeta.resolve(front)
    size_key = str(front.get("size") or "16x9").strip().lower()
    size = specs.SLIDE_SIZES_IN.get(size_key, specs.SLIDE_SIZES_IN["16x9"])
    base = template.base if template and template.base and template.base.suffix in (".pptx", ".potx") else None
    slides = split_slides(body)
    if not slides:
        raise ValueError("deck spec produced no slides — is the body empty?")
    builder = DeckBuilder(meta, size, base, spec_path.resolve().parent)
    builder.build(slides)
    builder.finish(out)
    return {
        "output": str(out),
        "template": template.name if template else None,
        "slides": len(slides),
        "layouts": [infer_layout(slide, index) for index, slide in enumerate(slides)],
        "unresolved_fields": missing,
        "warnings": builder.warnings,
    }


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(prog="build pptx", description="Markdown deck spec -> .pptx")
    parser.add_argument("spec", type=Path)
    parser.add_argument("-o", "--output", type=Path, required=True)
    parser.add_argument("--template", default="")
    parser.add_argument("--values", type=Path, default=None)
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args(argv)
    try:
        report = build(args.spec, args.output, template_name=args.template, values_path=args.values)
    except (ValueError, OSError, KeyError) as exc:
        print(f"pptx build failed: {exc}", file=sys.stderr)
        return 1
    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        print(f"wrote {report['output']}  ({report['slides']} slides: {', '.join(report['layouts'])})")
        for warning in report["warnings"]:
            print(f"  warning: {warning}")
        if report["unresolved_fields"]:
            print("  unresolved fields: " + ", ".join(report["unresolved_fields"]))
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
