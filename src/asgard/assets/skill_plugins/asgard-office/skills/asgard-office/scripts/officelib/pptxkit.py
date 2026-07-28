"""python-pptx helpers for the things PowerPoint stores as XML, not as API.

Bullets, slide backgrounds, and text-box insets have no Python surface in
python-pptx. They are three lines of OOXML each; open-coding them per call site
is how a deck ends up with literal bullet glyphs and double markers.
"""

from __future__ import annotations

from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# A rough advance-width per character, as a multiple of font size, per script.
# Used only to predict overflow before anyone opens the deck — see verify.py.
LATIN_ADVANCE = 0.52
CJK_ADVANCE = 1.0


def _ppr(paragraph):
    return paragraph._p.get_or_add_pPr()


def set_bullet(paragraph, char: str = "•", color: str | None = None, size_pct: int = 90) -> None:
    """Give a paragraph a real bullet. Never type the glyph into the text."""
    from pptx.oxml.ns import nsmap  # noqa: F401  (namespace registration side effect)
    from lxml import etree

    properties = _ppr(paragraph)
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        existing = properties.find(qn(tag))
        if existing is not None:
            properties.remove(existing)
    if color:
        fill = etree.SubElement(properties, qn("a:buClr"))
        srgb = etree.SubElement(fill, qn("a:srgbClr"))
        srgb.set("val", color)
    size = etree.SubElement(properties, qn("a:buSzPct"))
    size.set("val", str(int(size_pct * 1000)))
    node = etree.SubElement(properties, qn("a:buChar"))
    node.set("char", char)


def set_auto_number(paragraph, scheme: str = "arabicPeriod") -> None:
    from lxml import etree

    properties = _ppr(paragraph)
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        existing = properties.find(qn(tag))
        if existing is not None:
            properties.remove(existing)
    node = etree.SubElement(properties, qn("a:buAutoNum"))
    node.set("type", scheme)


def set_no_bullet(paragraph) -> None:
    from lxml import etree

    properties = _ppr(paragraph)
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        existing = properties.find(qn(tag))
        if existing is not None:
            properties.remove(existing)
    etree.SubElement(properties, qn("a:buNone"))


def set_indent(paragraph, level: int, *, hanging_pt: float = 16.0) -> None:
    properties = _ppr(paragraph)
    properties.set("marL", str(int(Pt(hanging_pt * (level + 1)))))
    properties.set("indent", str(int(-Pt(hanging_pt))))


def set_space(paragraph, *, before: float = 0.0, after: float = 6.0) -> None:
    from lxml import etree

    properties = _ppr(paragraph)
    for tag, value in (("a:spcBef", before), ("a:spcAft", after)):
        existing = properties.find(qn(tag))
        if existing is not None:
            properties.remove(existing)
        if value:
            holder = etree.SubElement(properties, qn(tag))
            points = etree.SubElement(holder, qn("a:spcPts"))
            points.set("val", str(int(value * 100)))


def set_background(slide, color: str) -> None:
    from pptx.dml.color import RGBColor

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)


def zero_inset(text_frame) -> None:
    """Text boxes carry a default inset; alignment against a shape needs it gone."""
    text_frame.margin_left = Emu(0)
    text_frame.margin_right = Emu(0)
    text_frame.margin_top = Emu(0)
    text_frame.margin_bottom = Emu(0)


def text_width_ratio(text: str) -> float:
    """Advance width of a string in units of font size (CJK glyphs are ~2x Latin)."""
    total = 0.0
    for char in text:
        total += CJK_ADVANCE if ord(char) > 0x2E7F else LATIN_ADVANCE
    return total


def estimate_lines(text: str, *, font_pt: float, box_width_in: float) -> int:
    """How many rendered lines a string needs in a box. Deliberately conservative."""
    if not text.strip():
        return 1
    capacity = max(1.0, (box_width_in * 72.0) / font_pt)
    return max(1, int(text_width_ratio(text) / capacity) + (1 if text_width_ratio(text) % capacity else 0))
