#!/usr/bin/env python3
"""The delivery gate: every defect a machine can prove without opening Word.

    asgard skills run asgard-office -- verify FILE [--strict] [--json]

Three bands, and the distinction is the whole point:

  error   the file is broken or ships something the author did not mean —
          a dangling relationship, a leftover {{placeholder}}, a #REF!.
  warn    a defect a reader will see — text past its box, a heading level
          skipped, a formula that cannot evaluate where it is going to be opened.
  info    something true that needs a human, not a fix — no cached formula
          values yet, tracked changes present, a scanned PDF.

`--strict` fails on warn as well. What this gate cannot see is layout: only a
renderer knows where a line actually broke. `render` is that gate, and it is
external on purpose — see references/qa.md.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
import zipfile
from dataclasses import asdict, dataclass
from pathlib import Path
from xml.etree import ElementTree

sys.path.insert(0, str(Path(__file__).resolve().parent))

# Leftovers that mean a draft shipped as a deliverable.
PLACEHOLDER_PROSE = re.compile(
    r"\blorem\b|\bipsum\b|\bTODO\b|\bTBD\b|\bFIXME\b|\[insert[^\]]*\]|\bXXX+\b|<제목>|<이름>|여기에 입력",
    re.I,
)
UNRESOLVED = re.compile(r"\{\{[^}]{0,80}\}\}")

# openpyxl writes formulas into the XML verbatim, and Excel stores every
# post-2007 name with an `_xlfn.` prefix its UI hides. Written bare, each of
# these lands in the delivered file as #NAME?.
NEEDS_XLFN = ("TEXTJOIN", "CONCAT", "IFS", "SWITCH", "MAXIFS", "MINIFS", "XOR", "IFNA")
# Dynamic arrays need spill metadata a library-written file does not carry, so
# only the top-left cell of the range ever gets a value.
SPILLS = ("XLOOKUP", "XMATCH", "FILTER", "SORT", "SORTBY", "UNIQUE", "SEQUENCE", "RANDARRAY", "TOCOL", "TOROW")
ERROR_LITERALS = ("#REF!", "#NAME?", "#VALUE!", "#DIV/0!", "#N/A", "#NULL!", "#NUM!")

OOXML_KINDS = {".docx": "docx", ".dotx": "docx", ".pptx": "pptx", ".potx": "pptx", ".xlsx": "xlsx", ".xlsm": "xlsx"}


@dataclass
class Finding:
    level: str  # error | warn | info
    code: str
    where: str
    message: str
    fix: str = ""


from officelib.color import contrast  # noqa: E402

# WCAG floors: 4.5:1 for body text, 3.0:1 once the text is large (>= 18pt, or
# >= 14pt bold). Below the large-text floor nothing is readable at any size.
BODY_CONTRAST = 4.5
LARGE_CONTRAST = 3.0
LARGE_PT = 18.0


# ------------------------------------------------------------------- package


def check_package(path: Path) -> list[Finding]:
    """Zip and OPC integrity — the failures that make a file simply not open."""
    out: list[Finding] = []
    try:
        archive = zipfile.ZipFile(path)
    except zipfile.BadZipFile:
        return [Finding("error", "not-a-package", path.name, "not a readable OOXML package (zip)", "rebuild the file")]
    with archive:
        names = set(archive.namelist())
        for name in names:
            if name.startswith("/") or ".." in Path(name).parts:
                out.append(
                    Finding("error", "unsafe-path", name, "package entry escapes the archive root", "rebuild the file")
                )
        if "[Content_Types].xml" not in names:
            out.append(
                Finding("error", "no-content-types", path.name, "[Content_Types].xml is missing", "rebuild the file")
            )
        for name in sorted(names):
            if not name.endswith((".xml", ".rels")):
                continue
            try:
                ElementTree.fromstring(archive.read(name))
            except ElementTree.ParseError as exc:
                out.append(Finding("error", "malformed-xml", name, f"XML will not parse: {exc}", "rebuild the part"))
        for name in sorted(name for name in names if name.endswith(".rels")):
            base = str(Path(name).parent.parent)
            try:
                root = ElementTree.fromstring(archive.read(name))
            except ElementTree.ParseError:
                continue
            for relationship in root:
                if relationship.get("TargetMode") == "External":
                    continue
                target = relationship.get("Target") or ""
                if target.startswith(("http:", "https:", "mailto:", "#")):
                    continue
                # Relationship targets are package-relative to the part's own folder;
                # a leading slash makes them package-absolute instead.
                resolved = target.lstrip("/") if target.startswith("/") else _normalise(base, target)
                if resolved not in names:
                    out.append(
                        Finding(
                            "error",
                            "dangling-relationship",
                            name,
                            f"relationship points at a part that is not in the package: {target}",
                            "remove the relationship or add the part",
                        )
                    )
    return out


def _normalise(base: str, target: str) -> str:
    parts: list[str] = [] if base in (".", "") else base.split("/")
    for piece in target.split("/"):
        if piece in ("", "."):
            continue
        if piece == "..":
            if parts:
                parts.pop()
        else:
            parts.append(piece)
    return "/".join(parts)


def _text_findings(kind: str, where: str, text: str) -> list[Finding]:
    out: list[Finding] = []
    for match in UNRESOLVED.finditer(text):
        out.append(
            Finding(
                "error",
                "unresolved-placeholder",
                where,
                f"template placeholder was never filled: {match.group(0)}",
                "supply the value, or delete the placeholder",
            )
        )
    match = PLACEHOLDER_PROSE.search(text)
    if match:
        out.append(
            Finding("warn", "draft-leftover", where, f"draft filler shipped in the text: {match.group(0)!r}",
                    "replace it with real content")
        )
    return out


# ---------------------------------------------------------------------- docx


def check_docx(path: Path) -> list[Finding]:
    from docx import Document

    out: list[Finding] = []
    document = Document(str(path))
    levels: list[int] = []
    headings = 0
    cjk = re.compile(r"[ᄀ-ᇿ぀-ヿ㄰-㆏가-힯一-鿿]")
    for index, paragraph in enumerate(document.paragraphs):
        text = paragraph.text
        out += _text_findings("docx", f"paragraph {index + 1}", text)
        style = (paragraph.style.name or "").lower()
        if style.startswith("heading"):
            headings += 1
            tail = style.split()[-1]
            level = int(tail) if tail.isdigit() else 1
            if not text.strip():
                out.append(
                    Finding("warn", "empty-heading", f"paragraph {index + 1}", "heading has no text", "remove it")
                )
            if levels and level > levels[-1] + 1:
                out.append(
                    Finding(
                        "warn",
                        "heading-jump",
                        f"paragraph {index + 1}",
                        f"heading level jumps {levels[-1]} -> {level}; the contents tree will read wrong",
                        "add the intermediate level or promote this one",
                    )
                )
            levels.append(level)
        if cjk.search(text):
            for run in paragraph.runs:
                if not cjk.search(run.text):
                    continue
                fonts = run._element.find(
                    "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}rPr"
                )
                east = None
                if fonts is not None:
                    node = fonts.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}rFonts")
                    east = node.get(
                        "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}eastAsia"
                    ) if node is not None else None
                if not east:
                    out.append(
                        Finding(
                            "warn",
                            "cjk-font-unset",
                            f"paragraph {index + 1}",
                            "CJK text with no w:eastAsia font — Word will substitute one, and the document "
                            "will not look the way it does here",
                            "set theme.font_cjk in the spec",
                        )
                    )
                break
    for table_index, table in enumerate(document.tables, start=1):
        for row in table.rows:
            for cell in row.cells:
                out += _text_findings("docx", f"table {table_index}", cell.text)

    body = ElementTree.tostring(document.element.body, encoding="unicode")
    if "TOC \\o" in body and not headings:
        out.append(
            Finding("warn", "empty-toc", path.name, "the document has a contents field but no headings",
                    "add headings, or drop the toc")
        )
    if "TOC \\o" in body:
        out.append(
            Finding("info", "toc-needs-update", path.name,
                    "the contents field is empty until a word processor updates it (select all, F9)",
                    "expected — Word populates it on open/print")
        )
    if "<w:ins " in body or "<w:del " in body:
        out.append(
            Finding("info", "tracked-changes", path.name, "the document carries tracked changes",
                    "accept or reject them before delivery if that was not intended")
        )
    return out


# ---------------------------------------------------------------------- pptx


def check_pptx(path: Path) -> list[Finding]:
    from pptx import Presentation

    from officelib import pptxkit

    out: list[Finding] = []
    presentation = Presentation(str(path))
    width = presentation.slide_width / 914400
    height = presentation.slide_height / 914400
    layouts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        where = f"slide {index}"
        boxes: list[tuple[float, float, float, float, str]] = []
        has_title = False
        has_visual = False
        background = _slide_background(slide) or "FFFFFF"
        for shape in slide.shapes:
            left = (shape.left or 0) / 914400
            top = (shape.top or 0) / 914400
            box_width = (shape.width or 0) / 914400
            box_height = (shape.height or 0) / 914400
            if shape.has_text_frame:
                text = shape.text_frame.text
                out += _text_findings("pptx", where, text)
                if text.strip() and not has_title:
                    has_title = True
                needed = 0.0
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text
                    if not line.strip():
                        needed += 8 / 72.0
                        continue
                    size = next((run.font.size.pt for run in paragraph.runs if run.font.size), 18.0)
                    lines = pptxkit.estimate_lines(line, font_pt=size, box_width_in=max(box_width, 0.4))
                    needed += lines * size * 1.22 / 72.0
                if box_height and needed > box_height * 1.28:
                    out.append(
                        Finding(
                            "warn",
                            "text-overflow",
                            f"{where} / {shape.name}",
                            f"text needs about {needed:.2f}in in a {box_height:.2f}in box",
                            "shorten the text, shrink the font, or grow the box",
                        )
                    )
                reported: set[str] = set()
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        colour = _run_colour(run)
                        if not colour or not run.text.strip() or colour in reported:
                            continue
                        size = run.font.size.pt if run.font.size else 18.0
                        large = size >= LARGE_PT or (size >= 14.0 and run.font.bold)
                        floor = LARGE_CONTRAST if large else BODY_CONTRAST
                        ratio = contrast(colour, background)
                        if ratio < floor:
                            reported.add(colour)
                            out.append(
                                Finding(
                                    "warn",
                                    "low-contrast",
                                    f"{where} / {shape.name}",
                                    f"#{colour} on #{background} is {ratio:.1f}:1 at {size:.0f}pt "
                                    f"(floor {floor:.1f}:1)",
                                    "darken the ink, lighten the ground, or make the text larger",
                                )
                            )
                if text.strip():
                    boxes.append((left, top, box_width, box_height, shape.name))
            else:
                has_visual = True
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            out += _text_findings("pptx", where, cell.text)
            if left < -0.05 or top < -0.05 or left + box_width > width + 0.05 or top + box_height > height + 0.05:
                out.append(
                    Finding(
                        "error",
                        "off-slide",
                        f"{where} / {shape.name}",
                        f"shape sits at ({left:.2f}, {top:.2f}) {box_width:.2f}x{box_height:.2f}in on a "
                        f"{width:.2f}x{height:.2f}in slide — PowerPoint writes it, it just is not visible",
                        "move it inside the canvas",
                    )
                )
        for first in range(len(boxes)):
            for second in range(first + 1, len(boxes)):
                if _overlap(boxes[first], boxes[second]) > 0.06:
                    out.append(
                        Finding(
                            "warn",
                            "overlap",
                            where,
                            f"{boxes[first][4]} and {boxes[second][4]} overlap by "
                            f"{_overlap(boxes[first], boxes[second]):.2f} sq in",
                            "separate them; text through text is the defect a reader sees first",
                        )
                    )
        if not has_title:
            out.append(Finding("warn", "untitled-slide", where, "the slide carries no text at all", "add a title"))
        if not has_visual and len(boxes) <= 1 and index > 1:
            out.append(
                Finding("info", "text-only-slide", where, "text-only slide — nothing for the eye to land on",
                        "a chart, an image, or a stat callout carries a point further than a bullet")
            )
        layouts.append(_layout_signature(slide))
    for index in range(2, len(layouts)):
        if layouts[index] == layouts[index - 1] == layouts[index - 2] and layouts[index] != "":
            out.append(
                Finding(
                    "info",
                    "layout-monotony",
                    f"slides {index - 1}-{index + 1}",
                    "three slides in a row with the same shape signature",
                    "vary the layout; a deck that never changes shape reads as filler",
                )
            )
            break
    return out


def _slide_background(slide) -> str:
    try:
        fill = slide.background.fill
        if fill.type is not None and fill.fore_color and fill.fore_color.rgb is not None:
            return str(fill.fore_color.rgb)
    except (AttributeError, TypeError, ValueError):
        pass
    return ""


def _run_colour(run) -> str:
    try:
        if run.font.color and run.font.color.rgb is not None:
            return str(run.font.color.rgb)
    except (AttributeError, TypeError, ValueError):
        pass
    return ""


def _overlap(first, second) -> float:
    left = max(first[0], second[0])
    top = max(first[1], second[1])
    right = min(first[0] + first[2], second[0] + second[2])
    bottom = min(first[1] + first[3], second[1] + second[3])
    return max(0.0, right - left) * max(0.0, bottom - top)


def _layout_signature(slide) -> str:
    return ",".join(sorted(f"{shape.shape_type}" for shape in slide.shapes))


# ---------------------------------------------------------------------- xlsx


def check_xlsx(path: Path) -> list[Finding]:
    import openpyxl

    out: list[Finding] = []
    formulas = openpyxl.load_workbook(str(path), data_only=False)
    try:
        cached = openpyxl.load_workbook(str(path), data_only=True)
    except Exception:
        cached = None
    total = uncached = 0
    for name in formulas.sheetnames:
        sheet = formulas[name]
        values = cached[name] if cached is not None else None
        by_row: dict[int, list[str]] = {}
        for row in sheet.iter_rows():
            for cell in row:
                text = cell.value
                if isinstance(text, str):
                    out += _text_findings("xlsx", f"{name}!{cell.coordinate}", text)
                if not isinstance(text, str) or not text.startswith("="):
                    continue
                total += 1
                by_row.setdefault(cell.row, []).append(cell.coordinate)
                upper = text.upper()
                for function in NEEDS_XLFN:
                    if re.search(rf"(?<![A-Z_.]){function}\s*\(", upper) and f"_XLFN.{function}" not in upper:
                        out.append(
                            Finding(
                                "warn",
                                "missing-xlfn",
                                f"{name}!{cell.coordinate}",
                                f"{function} is stored with an _xlfn. prefix; written bare it opens as #NAME?",
                                f"write =_xlfn.{function}(...)",
                            )
                        )
                for function in SPILLS:
                    if re.search(rf"(?<![A-Z_.]){function}\s*\(", upper):
                        out.append(
                            Finding(
                                "warn",
                                "spill-formula",
                                f"{name}!{cell.coordinate}",
                                f"{function} spills into neighbouring cells, and a library-written file carries "
                                "no spill metadata — only the anchor cell gets a value",
                                "use INDEX/MATCH, or sort and filter in the generator before writing cells",
                            )
                        )
                if re.search(r"\[\d+\]", text):
                    out.append(
                        Finding(
                            "error",
                            "external-link",
                            f"{name}!{cell.coordinate}",
                            "formula links to a separate workbook file that is not shipped with this one",
                            "inline the value, or deliver both files together",
                        )
                    )
                if values is not None:
                    cell_value = values[cell.coordinate].value
                    if cell_value is None:
                        uncached += 1
                    elif isinstance(cell_value, str) and cell_value.strip() in ERROR_LITERALS:
                        out.append(
                            Finding(
                                "error",
                                "formula-error",
                                f"{name}!{cell.coordinate}",
                                f"the cached result is {cell_value.strip()}",
                                "fix the formula; never ship a workbook with a live error in it",
                            )
                        )
        for row, coordinates in by_row.items():
            if len(coordinates) < 3:
                continue
            shapes = {re.sub(r"\d+", "#", str(sheet[coordinate].value).upper()) for coordinate in coordinates}
            if len(shapes) > 1:
                out.append(
                    Finding(
                        "warn",
                        "inconsistent-row",
                        f"{name} row {row}",
                        "formulas across this row are not the same shape — a single hand-edited cell mid-row is "
                        "the commonest silent modelling error",
                        "make the row uniform, or move the exception to its own labelled cell",
                    )
                )
    if total and uncached:
        out.append(
            Finding(
                "info",
                "no-cached-values",
                path.name,
                f"{uncached} of {total} formulas have no cached result, so pandas and any previewer read them "
                "as empty until a spreadsheet application recalculates the file",
                "open it once in Excel/LibreOffice, or run `render --recalc` where LibreOffice is available",
            )
        )
    return out


# ----------------------------------------------------------------------- pdf


def check_pdf(path: Path) -> list[Finding]:
    from pypdf import PdfReader

    out: list[Finding] = []
    reader = PdfReader(str(path))
    if reader.is_encrypted:
        out.append(Finding("info", "encrypted", path.name, "the PDF is encrypted", "supply the password to read it"))
        return out
    empty = 0
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        out += _text_findings("pdf", f"page {index}", text)
        if not text.strip():
            empty += 1
    if empty:
        out.append(
            Finding(
                "info",
                "no-text-layer",
                path.name,
                f"{empty} of {len(reader.pages)} pages carry no extractable text — scanned or image-only",
                "OCR is an external gate; nothing here can read those pages",
            )
        )
    fields = reader.get_fields() or {}
    if fields:
        unfilled = [name for name, field in fields.items() if not (field.get("/V") if hasattr(field, "get") else None)]
        out.append(
            Finding(
                "info" if not unfilled else "warn",
                "form-fields",
                path.name,
                f"{len(fields)} form fields, {len(unfilled)} still empty",
                "fill them, or confirm the blanks are intended",
            )
        )
    return out


# -------------------------------------------------------------------- driver


def verify(path: Path, lane: str = "") -> list[Finding]:
    if not path.is_file():
        raise ValueError(f"file not found: {path}")
    suffix = path.suffix.lower()
    kind = lane or OOXML_KINDS.get(suffix, "pdf" if suffix == ".pdf" else "")
    if not kind:
        raise ValueError(f"nothing to verify for {suffix} (docx, pptx, xlsx, pdf)")
    findings: list[Finding] = []
    if kind in ("docx", "pptx", "xlsx"):
        findings += check_package(path)
        if any(item.level == "error" and item.code in ("not-a-package", "no-content-types") for item in findings):
            return findings
    findings += {"docx": check_docx, "pptx": check_pptx, "xlsx": check_xlsx, "pdf": check_pdf}[kind](path)
    order = {"error": 0, "warn": 1, "info": 2}
    return sorted(findings, key=lambda item: (order[item.level], item.code, item.where))


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(prog="verify", description="Static delivery gate for office documents")
    parser.add_argument("file", type=Path)
    parser.add_argument("--lane", default="", choices=("", "docx", "pptx", "xlsx", "pdf"))
    parser.add_argument("--strict", action="store_true", help="fail on warnings as well as errors")
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args(argv)
    try:
        findings = verify(args.file, args.lane)
    except (ValueError, OSError, zipfile.BadZipFile) as exc:
        print(f"verify failed: {exc}", file=sys.stderr)
        return 2
    errors = sum(1 for item in findings if item.level == "error")
    warnings = sum(1 for item in findings if item.level == "warn")
    if args.json:
        print(
            json.dumps(
                {
                    "file": str(args.file),
                    "errors": errors,
                    "warnings": warnings,
                    "findings": [asdict(item) for item in findings],
                },
                ensure_ascii=False,
                indent=2,
            )
        )
    else:
        for item in findings:
            print(f"[{item.level}] {item.code} · {item.where}: {item.message}")
            if item.fix:
                print(f"         fix: {item.fix}")
        verdict = "clean" if not findings else f"{errors} error(s), {warnings} warning(s)"
        print(f"\n{args.file.name}: {verdict}")
    return 1 if errors or (args.strict and warnings) else 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
