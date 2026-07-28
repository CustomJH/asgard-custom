#!/usr/bin/env python3
"""Read any office document back out as Markdown or JSON.

    asgard skills run asgard-office -- read FILE [--format markdown|json] [--sheet NAME]

Reading is half of expert document work: you cannot edit a deck you have not
looked at, and you cannot check your own output without reading it back. The
Markdown form is for a person or a model to read; the JSON form keeps
coordinates, so an edit can be planned against it.

`.hwp` / `.hwpx` are not handled here — the bundled `hwpx` skill owns them.
"""

from __future__ import annotations

import argparse
import json
import sys
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

SUPPORTED = (".docx", ".dotx", ".pptx", ".potx", ".xlsx", ".xlsm", ".pdf", ".md", ".txt", ".csv", ".tsv")


def _cell_text(cell) -> str:
    return " ".join(paragraph.text.strip() for paragraph in cell.paragraphs if paragraph.text.strip())


def read_docx(path: Path) -> dict:
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(str(path))
    blocks: list[dict] = []
    # Walk the body in document order and wrap each child directly. Matching
    # `document.paragraphs` against the body by element identity looks equivalent
    # and is not: lxml hands out proxy objects, so the same node can arrive as a
    # different Python object and the paragraph silently disappears from the read.
    for child in document.element.body.iterchildren():
        tag = child.tag.split("}")[-1]
        if tag == "p":
            paragraph = Paragraph(child, document)
            if not paragraph.text.strip():
                continue
            style = (paragraph.style.name or "").lower()
            if style.startswith("heading"):
                level = int(style.split()[-1]) if style.split()[-1].isdigit() else 1
                blocks.append({"kind": "heading", "level": level, "text": paragraph.text.strip()})
            elif style.startswith("list"):
                tail = style.split()[-1]
                blocks.append(
                    {
                        "kind": "list_item",
                        "ordered": "number" in style,
                        "level": int(tail) - 1 if tail.isdigit() else 0,
                        "text": paragraph.text.strip(),
                    }
                )
            else:
                blocks.append({"kind": "para", "text": paragraph.text.strip()})
        elif tag == "tbl":
            table = Table(child, document)
            rows = [[_cell_text(cell) for cell in row.cells] for row in table.rows]
            blocks.append({"kind": "table", "rows": rows})
    core = document.core_properties
    return {
        "format": "docx",
        "path": str(path),
        "properties": {"title": core.title or "", "author": core.author or "", "subject": core.subject or ""},
        "blocks": blocks,
        "sections": len(document.sections),
    }


def read_pptx(path: Path) -> dict:
    from pptx import Presentation

    presentation = Presentation(str(path))
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines = [paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()]
                if lines:
                    shapes.append({"kind": "text", "name": shape.name, "lines": lines})
            elif shape.has_table:
                shapes.append(
                    {
                        "kind": "table",
                        "name": shape.name,
                        "rows": [[cell.text for cell in row.cells] for row in shape.table.rows],
                    }
                )
            elif shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
                shapes.append({"kind": "image", "name": shape.name})
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({"index": index, "shapes": shapes, "notes": notes})
    return {
        "format": "pptx",
        "path": str(path),
        "size_in": [presentation.slide_width / 914400, presentation.slide_height / 914400],
        "slides": slides,
    }


def read_xlsx(path: Path, sheet: str = "") -> dict:
    import openpyxl

    formulas = openpyxl.load_workbook(str(path), data_only=False)
    try:
        cached = openpyxl.load_workbook(str(path), data_only=True)
    except Exception:  # a workbook with no cached values at all
        cached = None
    sheets = []
    for name in formulas.sheetnames:
        if sheet and name != sheet:
            continue
        source = formulas[name]
        values = cached[name] if cached is not None else None
        rows = []
        for row in source.iter_rows():
            line = []
            for cell in row:
                entry: dict = {"ref": cell.coordinate, "value": cell.value}
                if isinstance(cell.value, str) and cell.value.startswith("="):
                    entry["formula"] = cell.value
                    entry["cached"] = values[cell.coordinate].value if values is not None else None
                line.append(entry)
            if any(entry["value"] is not None for entry in line):
                rows.append(line)
        sheets.append({"name": name, "dimensions": source.dimensions, "rows": rows})
    return {"format": "xlsx", "path": str(path), "sheets": sheets}


def read_pdf(path: Path) -> dict:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = [{"index": index, "text": (page.extract_text() or "").strip()} for index, page in enumerate(reader.pages, 1)]
    info = reader.metadata or {}
    return {
        "format": "pdf",
        "path": str(path),
        "properties": {key.lstrip("/"): str(value) for key, value in info.items()},
        "encrypted": reader.is_encrypted,
        "pages": pages,
    }


def read_text(path: Path) -> dict:
    return {"format": path.suffix.lstrip(".") or "txt", "path": str(path), "text": path.read_text(encoding="utf-8")}


def read(path: Path, sheet: str = "") -> dict:
    suffix = path.suffix.lower()
    if not path.is_file():
        raise ValueError(f"file not found: {path}")
    if suffix in (".docx", ".dotx"):
        return read_docx(path)
    if suffix in (".pptx", ".potx"):
        return read_pptx(path)
    if suffix in (".xlsx", ".xlsm"):
        return read_xlsx(path, sheet)
    if suffix == ".pdf":
        return read_pdf(path)
    if suffix in (".md", ".txt", ".csv", ".tsv"):
        return read_text(path)
    if suffix in (".hwp", ".hwpx"):
        raise ValueError("Korean documents belong to the hwpx skill: asgard skills run hwpx -- extract FILE")
    if suffix in (".doc", ".ppt", ".xls"):
        raise ValueError(f"legacy binary format {suffix} — convert it first (see `render --help` for the gate)")
    raise ValueError(f"unsupported format {suffix} (supported: {', '.join(SUPPORTED)})")


# ------------------------------------------------------------------- markdown


def _md_table(rows: list[list[str]]) -> list[str]:
    if not rows:
        return []
    width = max(len(row) for row in rows)
    padded = [list(row) + [""] * (width - len(row)) for row in rows]
    out = ["| " + " | ".join(cell.replace("|", "\\|").replace("\n", " ") for cell in padded[0]) + " |"]
    out.append("|" + "|".join(["---"] * width) + "|")
    for row in padded[1:]:
        out.append("| " + " | ".join(cell.replace("|", "\\|").replace("\n", " ") for cell in row) + " |")
    return out


def to_markdown(data: dict) -> str:
    lines: list[str] = []
    if data["format"] in ("docx", "dotx"):
        properties = data["properties"]
        if properties.get("title"):
            lines += [f"<!-- title: {properties['title']} -->"]
        for block in data["blocks"]:
            if block["kind"] == "heading":
                lines += ["", "#" * min(block["level"], 6) + " " + block["text"]]
            elif block["kind"] == "list_item":
                indent = "  " * int(block.get("level") or 0)
                lines.append(indent + ("1. " if block["ordered"] else "- ") + block["text"])
            elif block["kind"] == "table":
                lines += [""] + _md_table(block["rows"])
            else:
                lines += ["", block["text"]]
    elif data["format"] == "pptx":
        for slide in data["slides"]:
            lines += ["", f"<!-- slide: {slide['index']} -->"]
            for order, shape in enumerate(slide["shapes"]):
                if shape["kind"] == "table":
                    lines += _md_table(shape["rows"])
                elif shape["kind"] == "image":
                    lines.append(f"![{shape['name']}]()")
                else:
                    head = "## " if order == 0 else ""
                    lines.append(head + shape["lines"][0])
                    lines += [f"- {line}" for line in shape["lines"][1:]]
            if slide["notes"]:
                lines.append(f"<!-- notes: {slide['notes']} -->")
    elif data["format"] == "xlsx":
        for sheet in data["sheets"]:
            lines += ["", f"## {sheet['name']}", ""]
            rows = [
                [
                    (
                        f"{entry['formula']} -> {entry['cached']}"
                        if entry.get("formula")
                        else ("" if entry["value"] is None else str(entry["value"]))
                    )
                    for entry in row
                ]
                for row in sheet["rows"]
            ]
            lines += _md_table(rows)
    elif data["format"] == "pdf":
        for page in data["pages"]:
            lines += ["", f"<!-- page: {page['index']} -->", page["text"]]
    else:
        lines.append(data["text"])
    return "\n".join(lines).strip() + "\n"


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(prog="read", description="Office document -> Markdown or JSON")
    parser.add_argument("file", type=Path)
    parser.add_argument("--format", choices=("markdown", "json"), default="markdown")
    parser.add_argument("--sheet", default="", help="xlsx only: read one sheet")
    args = parser.parse_args(argv)
    try:
        data = read(args.file, args.sheet)
    except (ValueError, OSError, zipfile.BadZipFile) as exc:
        print(f"read failed: {exc}", file=sys.stderr)
        return 1
    print(json.dumps(data, ensure_ascii=False, indent=2, default=str) if args.format == "json" else to_markdown(data))
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
