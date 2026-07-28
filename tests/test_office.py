"""Asgard Office (Sága) — the document lane.

The lane scripts live under `assets/skill_plugins/` because they ship as a skill
bundle, not as package modules. They are still real code with real failure modes,
so the tests import them the way the entrypoint does: by putting the bundled
`scripts/` directory on `sys.path`.

What is pinned here:
  · the Markdown subset parser, including the directive-swallowing bug that made
    `<!-- notes: … -->` land inside a bullet
  · the placeholder engine and the template registry's scope shadowing
  · every one of the 23 genre skeletons actually builds
  · the delivery gate has teeth — each finding class is provoked on purpose
  · fill survives Word's run fragmentation and expands table-row loops
  · WCAG contrast is gamma-correct (the naive form called ordinary greys illegible)
"""

from __future__ import annotations

import json
import sys
import unittest
from pathlib import Path

_REPO = Path(__file__).resolve().parents[1]
_SKILL = _REPO / "src" / "asgard" / "assets" / "skill_plugins" / "asgard-office" / "skills" / "asgard-office"
_SCRIPTS = _SKILL / "scripts"
if str(_SCRIPTS) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS))

import build_docx  # noqa: E402
import build_pptx  # noqa: E402
import build_xlsx  # noqa: E402
import extract  # noqa: E402
import fill as fill_lane  # noqa: E402
import outline  # noqa: E402
import verify as verify_lane  # noqa: E402
from officelib import color, compose, mdblocks, specs, templates  # noqa: E402


def _codes(findings) -> set[str]:
    return {item.code for item in findings}


class TestMarkdownSubset(unittest.TestCase):
    def test_inline_runs_carry_style(self):
        runs = mdblocks.inline("plain **bold** *italic* `code` [link](http://x) ~~gone~~")
        by_text = {run.text: run for run in runs}
        self.assertTrue(by_text["bold"].bold)
        self.assertTrue(by_text["italic"].italic)
        self.assertTrue(by_text["code"].code)
        self.assertTrue(by_text["gone"].strike)
        self.assertEqual(by_text["link"].link, "http://x")

    def test_nested_lists_and_tables(self):
        blocks = mdblocks.parse("- one\n  - deep\n- two\n\n| a | b |\n|---|---|\n| 1 | 2 |\n")
        listing = next(block for block in blocks if block.kind == "list")
        self.assertEqual([item.level for item in mdblocks.normalise_levels(listing.items)], [0, 1, 0])
        table = next(block for block in blocks if block.kind == "table")
        self.assertEqual([mdblocks.plain(cell) for cell in table.header], ["a", "b"])
        self.assertEqual(len(table.rows), 1)

    def test_ordered_and_task_items(self):
        blocks = mdblocks.parse("1. first\n2. second\n")
        self.assertTrue(all(item.ordered for item in blocks[0].items))
        blocks = mdblocks.parse("- [x] done\n- [ ] open\n")
        self.assertEqual([item.checked for item in blocks[0].items], [True, False])

    def test_a_comment_ends_a_list_instead_of_joining_it(self):
        """Regression: a directive comment was absorbed as a lazy continuation line.

        The deck lane reads directives from comments, so swallowing one silently
        moved speaker notes into a bullet — invisible until the deck was opened.
        """
        blocks = mdblocks.parse("- item one\n<!-- notes: say this out loud -->\n")
        kinds = [block.kind for block in blocks]
        self.assertEqual(kinds, ["list", "comment"])
        self.assertEqual(mdblocks.plain(blocks[0].items[0].runs), "item one")
        self.assertEqual(blocks[1].text, "notes: say this out loud")

    def test_lazy_continuation_still_wraps(self):
        blocks = mdblocks.parse("- item one\n  wrapped onto the next line\n")
        self.assertEqual(mdblocks.plain(blocks[0].items[0].runs), "item one wrapped onto the next line")

    def test_rule_is_its_own_block(self):
        self.assertEqual([block.kind for block in mdblocks.parse("a\n\n---\n\nb\n")], ["para", "rule", "para"])


class TestSpecs(unittest.TestCase):
    def test_front_matter_split(self):
        meta, body = specs.split_frontmatter("---\ntitle: T\ntoc: true\n---\n# H\n")
        self.assertEqual(meta["title"], "T")
        self.assertTrue(meta["toc"])
        self.assertEqual(body.strip(), "# H")

    def test_unquoted_placeholder_in_front_matter_recovers(self):
        """`title: {{x}}` is a YAML flow mapping; the parser must not die on it."""
        meta, _ = specs.split_frontmatter("---\ntitle: {{title}}\nauthor: A\n---\nbody\n")
        self.assertEqual(meta["title"], "{{title}}")
        self.assertEqual(meta["author"], "A")

    def test_units_and_page_setup(self):
        self.assertAlmostEqual(specs.mm("1in", 0), 25.4)
        self.assertAlmostEqual(specs.mm("10", 0), 10.0)
        page = specs.PageSetup.resolve({"size": "letter", "orientation": "landscape", "margins": "10mm"})
        self.assertTrue(page.landscape)
        self.assertAlmostEqual(page.width_mm, 279.4)
        self.assertAlmostEqual(page.left_mm, 10.0)

    def test_theme_ignores_junk_colours(self):
        theme = specs.Theme.resolve({"primary": "not-a-colour", "accent": "#0a0B0c"})
        self.assertEqual(theme.primary, specs.Theme().primary)
        self.assertEqual(theme.accent, "0A0B0C")


class TestPlaceholderEngine(unittest.TestCase):
    def test_scalars_dotted_paths_and_missing(self):
        text, missing = templates.render("{{a}} / {{b.c}} / {{gone}}", {"a": 1, "b": {"c": "x"}})
        self.assertEqual(text, "1 / x / ")
        self.assertEqual(missing, ["gone"])

    def test_sections_repeat_and_invert(self):
        text, _ = templates.render("{{#rows}}[{{n}}]{{/rows}}{{^rows}}none{{/rows}}", {"rows": [{"n": 1}, {"n": 2}]})
        self.assertEqual(text, "[1][2]")
        text, _ = templates.render("{{#rows}}[{{n}}]{{/rows}}{{^rows}}none{{/rows}}", {"rows": []})
        self.assertEqual(text, "none")

    def test_unclosed_section_is_an_error(self):
        with self.assertRaises(ValueError):
            templates.render("{{#rows}}x", {"rows": [1]})

    def test_strict_render_refuses_gaps(self):
        with self.assertRaises(ValueError):
            templates.render("{{missing}}", {}, strict=True)

    def test_placeholders_listed_in_first_seen_order(self):
        self.assertEqual(templates.placeholders("{{b}} {{a}} {{b}} {{#c}}{{d}}{{/c}}"), ["b", "a", "c", "d"])


class TestTemplateRegistry(unittest.TestCase):
    def setUp(self):
        import tempfile

        self._temp = tempfile.TemporaryDirectory(prefix="asgard-office-test-")
        self.root = Path(self._temp.name)
        self.addCleanup(self._temp.cleanup)

    def _write(self, name: str, extra: str = "", kind: str = "docx") -> Path:
        directory = self.root / ".asgard" / "office" / "templates" / name
        directory.mkdir(parents=True)
        (directory / "template.toml").write_text(
            f'schema = 1\nname = "{name}"\nkind = "{kind}"\ntitle = "T"\n{extra}\n', encoding="utf-8"
        )
        (directory / "body.md").write_text('---\ntitle: "{{title}}"\n---\n\n# {{title}}\n', encoding="utf-8")
        return directory

    def test_discovery_and_scope(self):
        self._write("scoped")
        found, problems = templates.discover(self.root)
        self.assertEqual(problems, [])
        names = {item.name: item for item in found}
        self.assertEqual(names["scoped"].origin, "project")
        self.assertIn("report-ko", names, "bundled templates must still be visible")
        self.assertEqual(names["report-ko"].origin, "bundled")

    def test_project_shadows_bundled_by_name(self):
        self._write("report-ko")
        found, _ = templates.discover(self.root)
        shadowed = next(item for item in found if item.name == "report-ko")
        self.assertEqual(shadowed.origin, "project")

    def test_bad_manifest_is_reported_not_raised(self):
        directory = self.root / ".asgard" / "office" / "templates" / "broken"
        directory.mkdir(parents=True)
        (directory / "template.toml").write_text('schema = 1\nname = "broken"\nkind = "wat"\n', encoding="utf-8")
        found, problems = templates.discover(self.root)
        self.assertFalse(any(item.name == "broken" for item in found))
        self.assertTrue(any("broken" in problem for problem in problems))

    def test_field_schema_checks(self):
        self._write(
            "typed",
            'name_unused = 0\n[[fields]]\nkey = "n"\ntype = "number"\nrequired = true\n'
            '[[fields]]\nkey = "rows"\ntype = "table"\n',
        )
        template = templates.resolve("typed", self.root)
        self.assertEqual(templates.check(template, {"n": 3, "rows": [{"a": 1}]}), [])
        self.assertTrue(any("required" in problem for problem in templates.check(template, {})))
        self.assertTrue(any("must be a number" in problem for problem in templates.check(template, {"n": "x"})))
        self.assertTrue(
            any("row mappings" in problem for problem in templates.check(template, {"n": 1, "rows": ["x"]}))
        )
        self.assertTrue(any("not declared" in problem for problem in templates.check(template, {"n": 1, "typo": 1})))


class TestColor(unittest.TestCase):
    def test_wcag_luminance_is_gamma_corrected(self):
        """The naive weighting called #7B8794 on white 1.8:1 — an ordinary caption grey."""
        self.assertAlmostEqual(color.luminance("FFFFFF"), 1.0, places=6)
        self.assertAlmostEqual(color.luminance("000000"), 0.0, places=6)
        self.assertAlmostEqual(color.contrast("000000", "FFFFFF"), 21.0, places=2)
        # The naive form scored this 1.8:1; gamma-correct it is 3.66:1 — still under
        # the 4.5:1 body floor, which is why the default muted grey was darkened.
        self.assertAlmostEqual(color.contrast("7B8794", "FFFFFF"), 3.66, places=2)
        self.assertGreaterEqual(color.contrast(specs.Theme().muted, "FFFFFF"), 4.5)

    def test_readable_accent_keeps_going_until_it_clears_the_floor(self):
        fixed = color.readable_accent("F96167", "FFFFFF", floor=4.5)
        self.assertGreaterEqual(color.contrast(fixed, "FFFFFF"), 4.5)
        self.assertNotEqual(fixed, "000000", "it should darken the hue, not discard it")

    def test_ink_follows_ground(self):
        self.assertEqual(color.readable_ink("10151F"), "FFFFFF")
        self.assertEqual(color.readable_ink("FFFFFF"), "1F2933")


class TestDocxLane(unittest.TestCase):
    def setUp(self):
        import tempfile

        self._temp = tempfile.TemporaryDirectory(prefix="asgard-office-docx-")
        self.root = Path(self._temp.name)
        self.addCleanup(self._temp.cleanup)

    def _build(self, spec: str, name: str = "out.docx") -> tuple[Path, dict]:
        source = self.root / "spec.md"
        source.write_text(spec, encoding="utf-8")
        out = self.root / name
        return out, build_docx.build(source, out, root=self.root)

    def test_full_document_round_trips(self):
        out, report = self._build(
            "---\ntitle: T\nsubtitle: S\nauthor: A\ntoc: true\nnumber_headings: true\n"
            'theme: {font_cjk: "맑은 고딕"}\n---\n\n'
            "# 요약\n\n본문 **강조**.\n\n- 하나\n  - 둘\n\n| a | b |\n|---|---|\n| 1 | 2 |\n\n"
            "> 인용\n\n```\ncode\n```\n\n---\n\n# 다음\n"
        )
        self.assertTrue(out.is_file())
        self.assertEqual(report["headings"], 2)
        data = extract.read(out)
        headings = [block["text"] for block in data["blocks"] if block["kind"] == "heading"]
        self.assertEqual(headings, ["1. 요약", "2. 다음"])
        table = next(block for block in data["blocks"] if block["kind"] == "table")
        self.assertEqual(table["rows"][0], ["a", "b"])
        self.assertEqual(data["properties"]["title"], "T")

    def test_page_setup_reaches_the_section(self):
        from docx import Document

        out, _ = self._build("---\npage: {size: letter, margins: {left: 10mm}}\n---\n\ntext\n")
        section = Document(str(out)).sections[0]
        assert section.page_width is not None and section.left_margin is not None
        self.assertAlmostEqual(section.page_width.mm, 215.9, places=1)
        self.assertAlmostEqual(section.left_margin.mm, 10.0, places=1)

    def test_missing_image_warns_instead_of_failing(self):
        _, report = self._build("---\n---\n\n![alt](nope.png)\n")
        self.assertTrue(any("image not found" in warning for warning in report["warnings"]))

    def test_template_theme_merges_under_the_spec(self):
        directory = self.root / ".asgard" / "office" / "templates" / "themed"
        directory.mkdir(parents=True)
        (directory / "template.toml").write_text(
            'schema = 1\nname = "themed"\nkind = "docx"\n[theme]\nprimary = "112233"\naccent = "445566"\n'
            "[defaults]\ntoc = true\n",
            encoding="utf-8",
        )
        front, _, template, _ = compose.resolve_spec(
            self._spec('---\ntemplate: themed\ntheme: {accent: "AABBCC"}\n---\n\nbody\n'),
            "",
            None,
            self.root,
            ("docx", "md"),
        )
        assert template is not None
        self.assertEqual(template.name, "themed")
        self.assertEqual(front["theme"]["primary"], "112233", "template supplies what the spec omits")
        self.assertEqual(front["theme"]["accent"], "AABBCC", "the spec wins key by key")
        self.assertTrue(front["toc"], "template defaults apply")

    def _spec(self, text: str) -> Path:
        source = self.root / "themed-spec.md"
        source.write_text(text, encoding="utf-8")
        return source


class TestPptxLane(unittest.TestCase):
    def setUp(self):
        import tempfile

        self._temp = tempfile.TemporaryDirectory(prefix="asgard-office-pptx-")
        self.root = Path(self._temp.name)
        self.addCleanup(self._temp.cleanup)

    DECK = (
        "---\ntitle: D\nauthor: A\nsize: 16x9\nfooter: F\n---\n\n"
        "# D\nsubtitle line\n<!-- notes: speak this -->\n\n---\n\n"
        "## Numbers\n<!-- layout: stat -->\n- 99% :: uptime\n- 7% :: over\n\n---\n\n"
        "## Points\n- one\n- two\n\n---\n\n"
        "## Split\n<!-- layout: two-col -->\n- left\n- |||\n- right\n\n---\n\n"
        "## Grid\n| a | b |\n|---|---|\n| 1 | 2 |\n\n---\n\n"
        "## Close\n<!-- layout: section -->\n"
    )

    def _build(self) -> tuple[Path, dict]:
        source = self.root / "deck.md"
        source.write_text(self.DECK, encoding="utf-8")
        out = self.root / "deck.pptx"
        return out, build_pptx.build(source, out, root=self.root)

    def test_layouts_are_inferred(self):
        _, report = self._build()
        self.assertEqual(report["layouts"], ["title", "stat", "bullets", "two-col", "table", "section"])
        self.assertEqual(report["slides"], 6)

    def test_notes_land_in_the_notes_pane(self):
        out, _ = self._build()
        data = extract.read(out)
        self.assertEqual(data["slides"][0]["notes"], "speak this")
        self.assertEqual(len(data["slides"]), 6)

    def test_two_column_split_marker(self):
        out, _ = self._build()
        data = extract.read(out)
        lefts = [shape for shape in data["slides"][3]["shapes"] if shape["kind"] == "text"]
        joined = " ".join(line for shape in lefts for line in shape["lines"])
        self.assertIn("left", joined)
        self.assertIn("right", joined)
        self.assertNotIn("|||", joined, "the split marker must not survive into the slide")

    def test_unknown_layout_is_refused(self):
        source = self.root / "bad.md"
        source.write_text("---\ntitle: X\n---\n\n## S\n<!-- layout: carousel -->\n- a\n", encoding="utf-8")
        with self.assertRaises(ValueError):
            build_pptx.build(source, self.root / "bad.pptx", root=self.root)

    def test_generated_deck_passes_its_own_gate(self):
        """The builder must not ship defects the gate would refuse from anyone else."""
        out, _ = self._build()
        findings = verify_lane.verify(out)
        offending = [item for item in findings if item.level in ("error", "warn")]
        self.assertEqual(offending, [], f"builder produced gate findings: {[item.code for item in offending]}")


class TestXlsxLane(unittest.TestCase):
    def setUp(self):
        import tempfile

        self._temp = tempfile.TemporaryDirectory(prefix="asgard-office-xlsx-")
        self.root = Path(self._temp.name)
        self.addCleanup(self._temp.cleanup)

    def _build(self, spec: str) -> tuple[Path, dict]:
        source = self.root / "sheet.yaml"
        source.write_text(spec, encoding="utf-8")
        out = self.root / "book.xlsx"
        return out, build_xlsx.build(source, out, root=self.root)

    SPEC = (
        "title: M\nsheets:\n"
        "  - name: Assumptions\n"
        '    columns:\n      - {header: Item, width: 20}\n      - {header: Value, format: "0.0%"}\n'
        "    rows:\n      - [growth, 0.15]\n"
        "    inputs: [B2]\n"
        "    notes: {B2: from the plan}\n"
        "  - name: Model\n"
        "    columns: [{header: P}, {header: V}]\n"
        "    rows: [[Q1, 100]]\n"
        '    cells: {C2: "=B2*Assumptions!$B$2"}\n'
    )

    def test_formulas_are_written_as_formulas(self):
        import openpyxl

        out, report = self._build(self.SPEC)
        self.assertEqual(report["formula_cells"], 1)
        sheet = openpyxl.load_workbook(str(out))["Model"]
        self.assertEqual(sheet["C2"].value, "=B2*Assumptions!$B$2")

    def test_input_cells_get_the_model_convention(self):
        import openpyxl

        out, _ = self._build(self.SPEC)
        cell = openpyxl.load_workbook(str(out))["Assumptions"]["B2"]
        self.assertEqual(cell.font.color.rgb[-6:], build_xlsx.INPUT_FONT)
        self.assertEqual(cell.number_format, "0.0%")
        self.assertIsNotNone(cell.comment)

    def test_csv_rows_are_coerced(self):
        (self.root / "data.csv").write_text("a,1\nb,2\n", encoding="utf-8")
        out, _ = self._build("sheets:\n  - name: S\n    columns: [{header: K}, {header: V}]\n    rows: data.csv\n")
        import openpyxl

        sheet = openpyxl.load_workbook(str(out))["S"]
        self.assertEqual(sheet["B2"].value, 1, "a numeric CSV field must not land as text")

    def test_empty_spec_is_refused(self):
        with self.assertRaises(ValueError):
            self._build("title: nothing\n")


class TestDeliveryGate(unittest.TestCase):
    """Each finding class is provoked deliberately — a gate nobody can fail is not a gate."""

    def setUp(self):
        import tempfile

        self._temp = tempfile.TemporaryDirectory(prefix="asgard-office-gate-")
        self.root = Path(self._temp.name)
        self.addCleanup(self._temp.cleanup)

    def test_docx_unresolved_placeholder_and_heading_jump(self):
        from docx import Document

        document = Document()
        document.add_paragraph("Level one", style="Heading 1")
        document.add_paragraph("Level three", style="Heading 3")
        document.add_paragraph("Dear {{client}}, TODO finish this")
        out = self.root / "bad.docx"
        document.save(str(out))
        codes = _codes(verify_lane.verify(out))
        self.assertIn("unresolved-placeholder", codes)
        self.assertIn("heading-jump", codes)
        self.assertIn("draft-leftover", codes)

    def test_docx_cjk_without_east_asian_font(self):
        from docx import Document

        document = Document()
        document.add_paragraph("한국어 문장입니다")
        out = self.root / "cjk.docx"
        document.save(str(out))
        self.assertIn("cjk-font-unset", _codes(verify_lane.verify(out)))

    def test_pptx_overflow_offslide_and_contrast(self):
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        cramped = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.35))
        cramped.text_frame.word_wrap = True
        run = cramped.text_frame.paragraphs[0].add_run()
        run.text = "a sentence far too long for a three inch box at twenty points " * 3
        run.font.size = Pt(20)
        outside = slide.shapes.add_textbox(Inches(13.0), Inches(0.5), Inches(3), Inches(0.5))
        outside.text_frame.paragraphs[0].add_run().text = "off canvas"
        faint = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
        pale = faint.text_frame.paragraphs[0].add_run()
        pale.text = "barely there"
        pale.font.size = Pt(12)
        pale.font.color.rgb = RGBColor.from_string("EEEEEE")
        out = self.root / "bad.pptx"
        presentation.save(str(out))
        codes = _codes(verify_lane.verify(out))
        self.assertIn("text-overflow", codes)
        self.assertIn("off-slide", codes)
        self.assertIn("low-contrast", codes)

    def test_xlsx_formula_lint(self):
        import openpyxl

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet["A1"] = '=TEXTJOIN(",",TRUE,B1:B3)'
        sheet["A2"] = "=XLOOKUP(B1,C:C,D:D)"
        sheet["A3"] = "='[1]Other'!$B$2"
        out = self.root / "bad.xlsx"
        workbook.save(str(out))
        codes = _codes(verify_lane.verify(out))
        self.assertIn("missing-xlfn", codes)
        self.assertIn("spill-formula", codes)
        self.assertIn("external-link", codes)

    def test_xlsx_inconsistent_row(self):
        import openpyxl

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet["B2"], sheet["C2"], sheet["D2"] = "=B1*2", "=C1*2", "=SUM(A1:A9)"
        out = self.root / "row.xlsx"
        workbook.save(str(out))
        self.assertIn("inconsistent-row", _codes(verify_lane.verify(out)))

    def test_broken_package_is_an_error_not_a_crash(self):
        out = self.root / "junk.docx"
        out.write_bytes(b"not a zip at all")
        codes = _codes(verify_lane.verify(out))
        self.assertIn("not-a-package", codes)

    def test_unsupported_extension_is_refused(self):
        target = self.root / "note.rtf"
        target.write_text("x", encoding="utf-8")
        with self.assertRaises(ValueError):
            verify_lane.verify(target)


class TestFill(unittest.TestCase):
    def setUp(self):
        import tempfile

        self._temp = tempfile.TemporaryDirectory(prefix="asgard-office-fill-")
        self.root = Path(self._temp.name)
        self.addCleanup(self._temp.cleanup)

    def _form(self) -> Path:
        from docx import Document

        document = Document()
        paragraph = document.add_paragraph()
        for chunk in ["Dear ", "{{cli", "ent", "}},"]:  # Word fragments runs exactly like this
            paragraph.add_run(chunk)
        paragraph.runs[0].bold = True
        document.sections[0].footer.paragraphs[0].text = "{{company}}"
        table = document.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "item"
        table.rows[1].cells[0].text = "{{#lines}}{{name}}"
        table.rows[1].cells[1].text = "{{amount}}{{/lines}}"
        out = self.root / "form.docx"
        document.save(str(out))
        return out

    VALUES = {
        "client": "Nuriflex",
        "company": "Asgard",
        "lines": [{"name": "design", "amount": "1"}, {"name": "build", "amount": "2"}],
    }

    def test_fragmented_placeholder_is_filled_and_formatting_survives(self):
        from docx import Document

        form = self._form()
        out = self.root / "filled.docx"
        fill_lane.fill(form, self.VALUES, out)
        document = Document(str(out))
        self.assertEqual(document.paragraphs[0].text, "Dear Nuriflex,")
        self.assertTrue(document.paragraphs[0].runs[0].bold, "the untouched run keeps its formatting")
        self.assertEqual(document.sections[0].footer.paragraphs[0].text, "Asgard")

    def test_table_row_loop_repeats(self):
        from docx import Document

        form = self._form()
        out = self.root / "filled.docx"
        report = fill_lane.fill(form, self.VALUES, out)
        self.assertEqual(report["repeated_rows"], 2)
        table = Document(str(out)).tables[0]
        self.assertEqual([row.cells[0].text for row in table.rows], ["item", "design", "build"])
        self.assertEqual([row.cells[1].text for row in table.rows], ["", "1", "2"])

    def test_missing_values_are_reported_not_guessed(self):
        form = self._form()
        report = fill_lane.fill(form, {"client": "X"}, self.root / "partial.docx")
        self.assertIn("company", report["missing_fields"])
        self.assertTrue(report["unfilled_placeholders"])

    def test_scan_types_a_row_loop_as_a_table(self):
        typed = dict((name, kind) for name, kind, _ in fill_lane.scan_typed(self._form()))
        self.assertEqual(typed["lines"], "table")
        self.assertEqual(typed["client"], "text")
        self.assertNotIn("name", typed, "row-scope keys are not top-level fields")

    def test_filled_document_passes_the_gate(self):
        out = self.root / "filled.docx"
        fill_lane.fill(self._form(), self.VALUES, out)
        self.assertNotIn("unresolved-placeholder", _codes(verify_lane.verify(out)))


class TestGenres(unittest.TestCase):
    def setUp(self):
        import tempfile

        self._temp = tempfile.TemporaryDirectory(prefix="asgard-office-genre-")
        self.root = Path(self._temp.name)
        self.addCleanup(self._temp.cleanup)

    def test_every_genre_builds(self):
        """A skeleton that does not build is a trap laid for whoever tries it."""
        builders = {"docx": build_docx.build, "pptx": build_pptx.build, "xlsx": build_xlsx.build}
        values = self.root / "values.json"
        values.write_text(json.dumps({"title": "T", "author": "A", "date": "2026-07-28"}), encoding="utf-8")
        for row in outline.catalogue():
            with self.subTest(genre=row["genre"]):
                spec = self.root / f"{row['genre']}.md"
                spec.write_text(outline.skeleton(row["genre"]), encoding="utf-8")
                out = self.root / f"{row['genre']}.{row['kind']}"
                report = builders[row["kind"]](spec, out, values_path=values, root=self.root)
                self.assertTrue(out.is_file())
                self.assertEqual(report["unresolved_fields"], [])

    def test_korean_headings_differ_from_english(self):
        english = outline.skeleton("postmortem", language="en")
        korean = outline.skeleton("postmortem", language="ko")
        self.assertIn("# Timeline", english)
        self.assertIn("# 타임라인", korean)

    def test_guidance_comments_never_reach_the_document(self):
        spec = self.root / "adr.md"
        spec.write_text(outline.skeleton("adr"), encoding="utf-8")
        values = self.root / "v.json"
        values.write_text(json.dumps({"title": "T", "author": "A", "date": "d"}), encoding="utf-8")
        out = self.root / "adr.docx"
        build_docx.build(spec, out, values_path=values, root=self.root)
        rendered = extract.to_markdown(extract.read(out))
        self.assertNotIn("<!--", rendered.replace("<!-- title:", ""))
        self.assertNotIn("marketing", rendered, "section guidance must be dropped, not printed")

    def test_unknown_genre_names_the_alternatives(self):
        with self.assertRaises(ValueError) as caught:
            outline.skeleton("nope")
        self.assertIn("adr", str(caught.exception))


class TestBundledWiring(unittest.TestCase):
    def test_plugin_is_registered_and_runnable(self):
        from asgard.skill_registry import bundled_plugins

        plugin = bundled_plugins().get("asgard-office")
        assert plugin is not None, "the office plugin must validate against the plugin schema"
        self.assertEqual(plugin["skills"], ["asgard-office"])
        self.assertEqual(plugin["entrypoints"], {"asgard-office": "asgard_office.py"})
        self.assertIn("worker", plugin["routing"]["asgard-office"]["agents"])

    def test_every_reference_named_by_the_skill_is_reachable(self):
        """Named references must resolve through the resource accessor, not a file path.

        Client scaffolds (.claude/, .agents/) receive SKILL.md alone — the
        references stay in the installed package. A skill that tells the model to
        open `references/x.md` relative to itself sends it to a path that does not
        exist in Claude Code, Cursor, or Codex.
        """
        import re

        from asgard.skill_registry import show_skill_resource

        body = (_SKILL / "SKILL.md").read_text(encoding="utf-8")
        self.assertIn(
            "asgard skills show asgard-office --resource",
            body,
            "SKILL.md must name the accessor, since a relative path fails outside the package",
        )
        named = set(re.findall(r"references/([a-z0-9\-]+\.md)", body))
        self.assertTrue(named)
        for name in sorted(named):
            with self.subTest(reference=name):
                self.assertTrue((_SKILL / "references" / name).is_file(), f"missing reference: {name}")
                self.assertTrue(show_skill_resource(str(_REPO), "asgard-office", f"references/{name}").strip())

    def test_reaches_every_mode(self):
        """Native, Claude Code, Cursor, and Codex must all be able to see and run it.

        The three client modes are scaffolded from `client_skill_bodies`, and the
        native loop resolves through `resolve_skills`. A skill that lands in one
        surface and not the others is the parity failure this pins: the user asks
        for a report in Cursor and nothing happens, with no error to read.
        """
        from asgard.skill_registry import _compatible_agents, client_skill_bodies, resolve_skills

        # Client scaffolds (.claude/skills, .agents/skills) are built from this.
        # Default assignment is worker alone — every role carrying it by default
        # would put the whole document contract into four discovery contexts.
        self.assertIn("asgard-office", {name for name, _ in client_skill_bodies("worker", str(_REPO))})
        for agent in ("freyja", "thor", "mimir"):
            with self.subTest(agent=agent):
                self.assertIn(
                    agent,
                    _compatible_agents("asgard-office"),
                    "a design, backend, or walkthrough role must be able to opt in via `skills assign`",
                )
        self.assertNotIn(
            "eitri",
            _compatible_agents("asgard-office"),
            "build and release is not a document role",
        )

        # The native loop and the Codex/Cursor router both go through this.
        for task in ("제안서를 워드 문서로 만들어줘", "make a board deck", "장애 리뷰 보고서를 docx 로"):
            with self.subTest(task=task):
                self.assertIn("asgard-office", [name for name, _ in resolve_skills(str(_REPO), task, "worker")])
        self.assertNotIn(
            "asgard-office",
            [name for name, _ in resolve_skills(str(_REPO), "fix the flaky login test", "worker")],
            "a coding task must not drag the document lane in",
        )

    def test_bundled_templates_render(self):
        import tempfile

        with tempfile.TemporaryDirectory(prefix="asgard-office-bundled-") as temp:
            root = Path(temp)
            builders = {"docx": build_docx.build, "pptx": build_pptx.build, "xlsx": build_xlsx.build}
            for template in templates.discover(root)[0]:
                if template.origin != "bundled":
                    continue
                with self.subTest(template=template.name):
                    assert template.body is not None, "a bundled template ships a skeleton"
                    assert template.example is not None, "a bundled template ships a worked example"
                    out = root / f"{template.name}.{template.kind}"
                    report = builders[template.kind](
                        template.body, out, template_name=template.name, values_path=template.example, root=root
                    )
                    self.assertEqual(report["unresolved_fields"], [])
                    errors = [item for item in verify_lane.verify(out) if item.level == "error"]
                    self.assertEqual(errors, [], f"{template.name}: {[item.code for item in errors]}")


if __name__ == "__main__":
    unittest.main()
